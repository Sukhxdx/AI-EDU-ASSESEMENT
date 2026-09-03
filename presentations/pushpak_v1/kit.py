"""Pushpak V1 design-review visual kit — dark aerospace PDR language."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# --- palette (professional UAV design-review) ---
BG = "0A1220"
BG2 = "0C1728"
CARD = "121C30"
CARD2 = "172338"
STROKE = "243450"
CYAN = "00D4E8"
CYAN2 = "5AE8F4"
CYAN_DIM = "0D3C48"
AMBER = "F5B942"
ORANGE = "FF7A45"
TEAL = "2EE6A6"
RED = "FF5C6A"
WHITE = "FFFFFF"
OFF = "D7DFEA"
MUTED = "8A97AD"
DIM = "5C6B82"

FONT = "Inter"
MONO = "JetBrains Mono"
SW, SH = 13.333333, 7.5


def rgb(hx: str) -> RGBColor:
    hx = hx.lstrip("#")
    return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))


def _set_run(run, text, size, color, bold=False, font=FONT, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    run.font.name = font
    # disable theme latin so custom font sticks
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", font)


def set_bg(slide, hx=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hx)


def no_line(shape):
    shape.line.fill.background()


def shape_fill(shape, hx):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hx)


def shape_line(shape, hx, pt=1.0):
    shape.line.color.rgb = rgb(hx)
    shape.line.width = Pt(pt)


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape_fill(sh, fill)
    if line:
        shape_line(sh, line, line_w)
    else:
        no_line(sh)
    sh.shadow.inherit = False
    return sh


def rrect(slide, x, y, w, h, fill, line=None, line_w=1.0, adj=0.08):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape_fill(sh, fill)
    if line:
        shape_line(sh, line, line_w)
    else:
        no_line(sh)
    try:
        sh.adjustments[0] = adj
    except Exception:
        pass
    sh.shadow.inherit = False
    return sh


def oval(slide, x, y, w, h, fill, line=None, line_w=1.25):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape_fill(sh, fill)
    if line:
        shape_line(sh, line, line_w)
    else:
        no_line(sh)
    sh.shadow.inherit = False
    return sh


def chevron(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape_fill(sh, fill)
    no_line(sh)
    sh.shadow.inherit = False
    return sh


def textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=14,
    color=WHITE,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic=False,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    _set_run(p.add_run(), text, size, color, bold, font, italic)
    # tighten margins
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb


def multiline(
    slide,
    x,
    y,
    w,
    h,
    lines,
    size=12,
    color=OFF,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    spacing=1.05,
    anchor=MSO_ANCHOR.TOP,
):
    """lines: list of str or (str, kwargs) """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        p.line_spacing = spacing
        if isinstance(item, tuple):
            txt, kw = item
            _set_run(
                p.add_run(),
                txt,
                kw.get("size", size),
                kw.get("color", color),
                kw.get("bold", bold),
                kw.get("font", font),
                kw.get("italic", False),
            )
        else:
            _set_run(p.add_run(), item, size, color, bold, font)
    return tb


def notes(slide, body: str):
    ns = slide.notes_slide
    ns.notes_text_frame.text = body


def header(slide, kicker: str, title: str, subtitle: str | None = None):
    rect(slide, 0, 0, SW, 0.055, CYAN)
    textbox(slide, 0.42, 0.16, 12.4, 0.28, kicker.upper(), 10, CYAN, True, MONO)
    textbox(slide, 0.40, 0.38, 12.5, 0.46, title, 24, WHITE, True, FONT)
    if subtitle:
        textbox(slide, 0.42, 0.82, 12.4, 0.28, subtitle, 12, MUTED, False, FONT)


def footer(slide, num: int, total: int = 40):
    rect(slide, 0, 7.38, SW, 0.12, BG2)
    rect(slide, 0, 7.38, SW, 0.015, STROKE)
    textbox(
        slide,
        0.42,
        7.39,
        9.2,
        0.10,
        "PUSHPAK V1   ·   TILT-ROTOR VTOL DEMONSTRATOR   ·   FINAL YEAR DESIGN REVIEW",
        8,
        DIM,
        False,
        MONO,
        PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE,
    )
    textbox(
        slide,
        10.55,
        7.39,
        2.35,
        0.10,
        f"{num:02d}  /  {total:02d}",
        9,
        CYAN,
        True,
        MONO,
        PP_ALIGN.RIGHT,
        MSO_ANCHOR.MIDDLE,
    )


def card(slide, x, y, w, h, accent=None, fill=CARD, line=STROKE):
    sh = rrect(slide, x, y, w, h, fill, line, 1.0, 0.06)
    if accent:
        rect(slide, x, y, w, 0.045, accent)
    return sh


def hud_corners(slide, x, y, w, h, arm=0.14, thick=0.018, color=CYAN):
    # four L-brackets
    rect(slide, x, y, arm, thick, color)
    rect(slide, x, y, thick, arm, color)
    rect(slide, x + w - arm, y, arm, thick, color)
    rect(slide, x + w - thick, y, thick, arm, color)
    rect(slide, x, y + h - thick, arm, thick, color)
    rect(slide, x, y + h - arm, thick, arm, color)
    rect(slide, x + w - arm, y + h - thick, arm, thick, color)
    rect(slide, x + w - thick, y + h - arm, thick, arm, color)


def kpi(slide, x, y, w, h, value, label, sub=None, accent=CYAN):
    card(slide, x, y, w, h, accent)
    textbox(slide, x + 0.12, y + 0.16, w - 0.24, 0.48, str(value), 22, WHITE, True, MONO)
    textbox(slide, x + 0.12, y + 0.62, w - 0.24, 0.28, label.upper(), 10, MUTED, True, FONT)
    if sub:
        textbox(slide, x + 0.12, y + h - 0.36, w - 0.24, 0.28, sub, 10, OFF, False, FONT)


def pill(slide, x, y, w, h, text, fill=CYAN_DIM, color=CYAN):
    rrect(slide, x, y, w, h, fill, None, 1, 0.5)
    textbox(slide, x, y, w, h, text, 10, color, True, MONO, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def bullet_block(slide, x, y, w, h, items, size=13, color=OFF):
    lines = []
    for it in items:
        lines.append(f"▸  {it}")
    return multiline(slide, x, y, w, h, lines, size, color)


def style_table(table, header=True, font_size=10, header_fill=CYAN_DIM, zebra=True):
    for ri, row in enumerate(table.rows):
        for ci, cell in enumerate(row.cells):
            cell.fill.solid()
            if ri == 0 and header:
                cell.fill.fore_color.rgb = rgb(header_fill)
                fg, b = CYAN, True
            else:
                cell.fill.fore_color.rgb = rgb(CARD if (not zebra or ri % 2 == 1) else CARD2)
                fg, b = OFF, False
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.size = Pt(font_size if ri else font_size)
                    run.font.bold = b or (ri == 0)
                    run.font.color.rgb = rgb(fg if ri else CYAN)
                    run.font.name = MONO if ci > 0 or ri == 0 else FONT
            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_right = Inches(0.06)
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # borders
    _set_table_borders(table, STROKE)


def _set_table_borders(table, hx):
    # tblPr / tcTcPr line
    tbl = table._tbl
    nsmap_a = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for cell in tbl.iter(qn("a:tc")):
        tcPr = cell.find(qn("a:tcPr"))
        if tcPr is None:
            tcPr = etree.SubElement(cell, qn("a:tcPr"))
        for edge in ("lnL", "lnR", "lnT", "lnB"):
            ln = tcPr.find(qn(f"a:{edge}"))
            if ln is not None:
                tcPr.remove(ln)
            ln = etree.SubElement(tcPr, qn(f"a:{edge}"))
            ln.set("w", "6350")  # 0.5 pt
            sf = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
            srgb.set("val", hx)


def add_table(slide, data, x, y, w, h, font_size=10):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
    style_table(table, font_size=font_size)
    return table


def picture(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def process_bar(slide, x, y, w, h, steps, active=None):
    n = len(steps)
    gap = 0.08
    cw = (w - gap * (n - 1)) / n
    for i, st in enumerate(steps):
        xx = x + i * (cw + gap)
        fill = CYAN_DIM if (active is None or i <= active) else CARD
        fg = CYAN if (active is None or i <= active) else MUTED
        ch = chevron(slide, xx, y, cw, h, fill)
        textbox(slide, xx + 0.06, y, cw - 0.18, h, st, 10, fg, True, FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return


def arrow_right(slide, x, y, w, h, fill=CYAN):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape_fill(sh, fill)
    no_line(sh)
    sh.shadow.inherit = False
    return sh
