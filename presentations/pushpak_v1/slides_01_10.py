"""Slides 01–10 — title through overall architecture."""

from pathlib import Path

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

import kit as k

A = Path(__file__).parent / "assets"


def _base(prs, kicker, title, num, notes, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    k.set_bg(s)
    k.header(s, kicker, title, subtitle)
    k.footer(s, num)
    k.notes(s, notes)
    return s


def s01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    k.set_bg(s)
    k.picture(s, A / "pushpak-hero.png", 4.55, 0, 8.783, 7.5)
    k.rect(s, 0, 0, 6.35, 7.5, k.BG)
    k.rect(s, 6.25, 0, 0.35, 7.5, k.BG)  # blend strip
    k.rect(s, 0, 0, 13.333, 0.055, k.CYAN)

    k.pill(s, 0.50, 1.15, 2.55, 0.32, "DESIGN REVIEW  ·  2026")
    k.textbox(s, 0.48, 1.62, 6.0, 0.28, "PUSHPAK  V1", 13, k.CYAN, True, k.MONO)
    k.textbox(s, 0.45, 1.92, 6.1, 1.35, "Tilt-Rotor VTOL\nUAV Demonstrator", 32, k.WHITE, True)
    k.multiline(
        s, 0.48, 3.45, 5.7, 0.7,
        ["Mission: demonstrate electrical-power reduction in forward",
         "flight through wing-borne lift vs pure rotor-borne flight."],
        13, k.OFF,
    )

    specs = [("4 ROTORS", "Front + rear tilt pairs"),
             ("4.5 kg", "Design-point MTOW"),
             ("1.0 kg", "Payload"),
             ("PIXHAWK", "ArduPilot + ESP32")]
    for i, (v, lab) in enumerate(specs):
        x = 0.50 + (i % 2) * 2.55
        y = 4.30 + (i // 2) * 0.95
        k.card(s, x, y, 2.42, 0.85, k.CYAN)
        k.textbox(s, x + 0.12, y + 0.10, 2.2, 0.36, v, 16, k.WHITE, True, k.MONO)
        k.textbox(s, x + 0.12, y + 0.46, 2.2, 0.28, lab.upper(), 9, k.MUTED, True)

    k.textbox(s, 0.50, 6.30, 5.7, 0.22, "PROJECT TEAM", 9, k.CYAN, True, k.MONO)
    k.multiline(
        s, 0.48, 6.50, 5.8, 0.72,
        ["Ishaan Bafna   ·   PRN 23070125017",
         "Aman Ravat     ·   PRN 23070125004",
         "Krit Verma     ·   PRN 230701250"],
        12, k.OFF, font=k.MONO,
    )
    k.footer(s, 1)
    k.notes(
        s,
        "Open on the aircraft, then the mission in one sentence: we are not building a product UAV — we are proving that wing-borne lift cuts electrical power versus a pure multirotor of the same mass.\n"
        "Introduce the team and PRNs. Flag MTOW 4.5 kg and 1 kg payload as the frozen design point.\n"
        "Transition: the next slide frames the whole system before we dive into the problem.",
    )


def s02_overview(prs):
    s = _base(
        prs, "01  /  PROJECT", "Project Overview", 2,
        "Walk the four identity cards left to right. Emphasize that Pushpak V1 is a demonstrator, not an operational ISR air vehicle.\n"
        "Call out the mission banner — every later requirement traces to power reduction via the wing.\n"
        "Close on the configuration: two tilt axes, 75° stop, Pixhawk/ArduPilot.",
        "RC tilt-rotor VTOL demonstrator  ·  four rotors  ·  fixed wing",
    )
    k.card(s, 0.42, 1.22, 12.48, 0.72, k.CYAN)
    k.textbox(
        s, 0.62, 1.35, 12.1, 0.48,
        "Prove that wing-borne lift reduces electrical power in forward flight relative to pure rotor-borne flight at the same MTOW.",
        15, k.WHITE, True,
    )
    items = [
        ("AIRCRAFT", "RC tilt-rotor VTOL\n4 rotors + fixed wing\nFront & rear tilt pairs"),
        ("DESIGN POINT", "MTOW  4.5 kg\nTarget  4.0–4.5 kg\nPayload  1.0 kg"),
        ("PROPULSION", "4 × MY4215 400KV\n16-inch class props\n6S Li-ion  ~20 Ah"),
        ("AVIONICS", "Pixhawk + ArduPilot\nESP32 companion\nPitot · GPS · mag encoders"),
        ("TILT", "2 × STS3215-C018\nOne servo per axis\nMax tilt  75°"),
        ("GEOMETRY", "Wing span  1300 mm\nFuselage  1200–1300 mm\nRotor span  1700–1900 mm"),
    ]
    for i, (h, body) in enumerate(items):
        x = 0.42 + (i % 3) * 4.16
        y = 2.12 + (i // 3) * 2.45
        k.card(s, x, y, 3.98, 2.28, k.CYAN)
        k.textbox(s, x + 0.18, y + 0.16, 3.6, 0.32, h, 12, k.CYAN, True, k.MONO)
        k.multiline(s, x + 0.12, y + 0.52, 3.7, 1.6, body.split("\n"), 15, k.OFF)


def s03_problem(prs):
    s = _base(
        prs, "02  /  PROBLEM", "Problem Statement", 3,
        "Multirotors hover well and waste energy in cruise because the rotors still carry 100% of weight.\n"
        "Fixed-wing aircraft cruise efficiently but cannot take off or land without a runway.\n"
        "Pushpak closes that gap at 4.5 kg so we can measure the power delta on one airframe.",
    )
    k.card(s, 0.42, 1.22, 6.05, 5.85)
    k.pill(s, 0.62, 1.42, 1.7, 0.30, "CONSTRAINT")
    k.textbox(s, 0.62, 1.88, 5.6, 0.9, "Rotor-borne flight\ncarries the full weight\nthrough the propellers.", 22, k.WHITE, True)
    k.bullet_block(s, 0.55, 3.35, 5.7, 3.3, [
        "Induced power scales with disk loading, not with airspeed.",
        "A 4.5 kg quadrotor still burns hundreds of watts just to stay aloft.",
        "Edgewise forward flight adds profile and parasite penalties.",
        "Endurance and range collapse before a useful payload mission is complete.",
        "Student-class VTOL demonstrators rarely isolate the wing-lift benefit with flight data.",
    ], 14)

    k.card(s, 6.70, 1.22, 6.20, 5.85, k.AMBER)
    k.pill(s, 6.90, 1.42, 1.9, 0.30, "RESPONSE", k.CYAN_DIM, k.AMBER)
    k.textbox(s, 6.90, 1.88, 5.8, 0.9, "Add a wing. Offload lift.\nMeasure the electrical\npower that disappears.", 22, k.WHITE, True)
    k.bullet_block(s, 6.82, 3.35, 5.9, 3.3, [
        "Wing target: 70% of weight at 20 m/s cruise.",
        "Rotors unload, then tilt to 75° to provide thrust.",
        "Same airframe flown wing-on vs wing-off (or rotor-borne equivalent).",
        "Success = a documented reduction in electrical power, not a paper polar.",
        "Demonstrator mass, cost and complexity stay inside a university build.",
    ], 14)


def s04_why(prs):
    s = _base(
        prs, "03  /  CONCEPT", "Why VTOL + Fixed Wing", 4,
        "Do not sell tilt-rotor as universally best — sell it as the right demonstrator: hover like a quad, cruise like a wing, one propulsion set.\n"
        "Tailsitters are simpler but ugly in hover handling. Tilt-wing is heavier. Dual-system (lift + cruise props) duplicates motors.\n"
        "Pushpak uses two coupled tilt axes so we can instrument both front and rear independently.",
    )
    data = [
        ["Attribute", "Multirotor", "Fixed wing", "Tilt-rotor VTOL"],
        ["Runway", "None", "Required", "None"],
        ["Hover", "Native", "No", "Native (4 rotors)"],
        ["Cruise lift", "Rotor-borne", "Wing-borne", "Wing + residual rotor"],
        ["Cruise power", "High", "Low", "Target: reduced"],
        ["Propulsion sets", "1", "1", "1 (reused)"],
        ["Transition", "N/A", "N/A", "Scheduled tilt + throttle"],
        ["Instrumentation", "Easy", "Easy", "Tilt angle + airspeed + power"],
        ["Fit to mission", "No (no wing benefit)", "No (no VTOL)", "Selected"],
    ]
    k.add_table(s, data, 0.42, 1.22, 12.48, 5.0, 11)
    k.card(s, 0.42, 6.38, 12.48, 0.85, k.CYAN)
    k.textbox(
        s, 0.62, 6.52, 12.1, 0.55,
        "Selection logic  ·  one motor set  ·  two tilt axes  ·  measurable wing-offload  ·  ArduPilot VTOL heritage",
        14, k.WHITE, True,
    )


def s05_mission(prs):
    s = _base(
        prs, "04  /  MISSION", "Mission Requirements", 5,
        "MR-1 is the only success criterion that cannot be traded. Everything else exists to make that measurement honest.\n"
        "Call out payload 1 kg as a mass stand-in — it forces a real MTOW, not a bare airframe.\n"
        "75° is a mechanical stop, not a cruise-optimal number; we will map power vs tilt in flight test.",
    )
    reqs = [
        ("MR-1", "PRIMARY", k.CYAN, "Quantify electrical-power reduction in wing-borne forward flight versus rotor-borne flight at the same mass."),
        ("MR-2", "VTOL", k.TEAL, "Vertical take-off and landing from an unprepared 10 × 10 m pad, no catapult, no net."),
        ("MR-3", "PAYLOAD", k.AMBER, "Carry 1.0 kg payload at design-point MTOW 4.5 kg (empty aircraft ≤ 3.5 kg)."),
        ("MR-4", "CRUISE", k.CYAN, "Sustain 20 m/s with the wing producing ~70% of lift (CL ≈ 0.395 on 0.32 m²)."),
        ("MR-5", "TRANSITION", k.TEAL, "Stable hover → 75° tilt transition using Pixhawk / ArduPilot with logged states."),
        ("MR-6", "DATA", k.AMBER, "Log airspeed, tilt, RPM, bus voltage/current, actuator current, and inertial rates."),
    ]
    for i, (rid, tag, col, txt) in enumerate(reqs):
        y = 1.20 + i * 0.95
        k.card(s, 0.42, y, 12.48, 0.88, col)
        k.textbox(s, 0.58, y + 0.22, 1.3, 0.44, rid, 18, col, True, k.MONO)
        k.pill(s, 1.95, y + 0.28, 1.55, 0.32, tag, k.CYAN_DIM, col)
        k.textbox(s, 3.70, y + 0.18, 8.9, 0.55, txt, 14, k.OFF)


def s06_objectives(prs):
    s = _base(
        prs, "05  /  OBJECTIVES", "Design Objectives", 6,
        "O1 is the experiment. O2–O6 are the vehicle that makes the experiment possible.\n"
        "If mass grows past 4.9 kg we descope structure or battery — we do not silently drop the 70% lift target.\n"
        "Safety objectives are equal priority for a student-flown 4.5 kg VTOL.",
    )
    objs = [
        ("O1", "Power experiment", "Wing-on vs wing-off (or rotor-borne) electrical-power comparison at matched mass and speed."),
        ("O2", "Thrust margin", "≥ 2.0× hover thrust (target 9 kgf static on 4.5 kg)."),
        ("O3", "Wing-offload", "NACA 4412, 0.32 m², AR 5.3, CL 0.395 at 20 m/s."),
        ("O4", "Tilt authority", "Front and rear axes independent, 0–75°, magnetic angle closed-loop."),
        ("O5", "Open-source FC", "ArduPilot VTOL stack on Pixhawk; ESP32 for high-rate tilt / power I/O."),
        ("O6", "Buildable", "University workshop: 3D-printed mounts, composite/foam wing, COTS propulsion."),
        ("O7", "Traceable mass", "Empty ≤ 3.5 kg; live mass budget from CAD + weigh-ins."),
        ("O8", "Safe envelope", "Mechanical tilt stops, current limits, geofence, RTL, dual-axis jam detect."),
    ]
    for i, (oid, title, body) in enumerate(objs):
        x = 0.42 + (i % 4) * 3.20
        y = 1.22 + (i // 4) * 2.90
        k.card(s, x, y, 3.05, 2.70, k.CYAN)
        k.textbox(s, x + 0.16, y + 0.16, 2.7, 0.32, oid, 14, k.CYAN, True, k.MONO)
        k.textbox(s, x + 0.16, y + 0.50, 2.7, 0.55, title, 16, k.WHITE, True)
        k.textbox(s, x + 0.12, y + 1.15, 2.78, 1.35, body, 13, k.OFF)


def s07_sysreq(prs):
    s = _base(
        prs, "06  /  REQUIREMENTS", "System Requirements", 7,
        "This is the shall-statement set we will verify. SR-P propulsion, SR-A airframe, SR-C control, SR-T test.\n"
        "Hover 1.125 kgf per motor is the physics floor; 2.25 kgf is the design target with margin.\n"
        "Verification column tells faculty we are not leaving this as a paper airplane.",
    )
    data = [
        ["ID", "Requirement", "Target", "Verification"],
        ["SR-P1", "Hover thrust per motor", "≥ 1.125 kgf  (design 2.25)", "Static thrust stand"],
        ["SR-P2", "Total static thrust", "≈ 9 kgf", "4-motor stand, 6S"],
        ["SR-P3", "Propulsion mass", "Motors 0.80 kg; ESC 0.25 kg", "Scale + CAD"],
        ["SR-A1", "Wing area / AR / airfoil", "0.32 m² · 5.3 · NACA 4412", "CAD + template"],
        ["SR-A2", "Cruise CL at 20 m/s", "≈ 0.395 (70% weight)", "CFD + pitot + load"],
        ["SR-A3", "Rotor / wing clearance", "50–75 mm min", "CAD + ground spin"],
        ["SR-C1", "Tilt range / axes", "0–75° · front + rear", "Encoder log"],
        ["SR-C2", "Flight stack", "Pixhawk + ArduPilot + ESP32", "HIL + taxi"],
        ["SR-T1", "Power comparison", "Wing-on vs wing-off", "Flight / trolley test"],
        ["SR-S1", "Empty mass", "≤ 3.5 kg", "Weigh-in"],
    ]
    k.add_table(s, data, 0.42, 1.18, 12.48, 5.95, 11)


def s08_config(prs):
    s = _base(
        prs, "07  /  CONFIGURATION", "Aircraft Configuration Selection", 8,
        "Four classical VTOL families. Dual-system wastes mass on a second propulsion set — fatal at 4.5 kg.\n"
        "Tilt-wing moves the entire wing; too much actuator work for this class.\n"
        "Tailsitter is the cheapest and the worst hover-to-cruise handling for a first student VTOL.\n"
        "Tilt-rotor with paired front/rear shafts is the selected architecture.",
    )
    opts = [
        ("A  TAILSITTER", False, "Single attitude change. Poor hover precision. Hard to instrument a clean wing-off comparison."),
        ("B  TILT-WING", False, "Wing + rotors tilt together. High hinge moment, complex structure, poor student-build risk."),
        ("C  DUAL-SYSTEM", False, "Lift rotors + cruise prop. Two motor sets. Mass and power electronics kill the 3.5 kg empty target."),
        ("D  TILT-ROTOR", True, "Four rotors, two tilt shafts. Hover as a quad. Cruise on the wing. One propulsion set. SELECTED."),
    ]
    for i, (title, sel, body) in enumerate(opts):
        x = 0.42 + i * 3.20
        k.card(s, x, 1.22, 3.05, 4.55, k.CYAN if sel else None, fill=k.CYAN_DIM if sel else k.CARD)
        if sel:
            k.pill(s, x + 0.16, 1.40, 1.45, 0.28, "SELECTED")
        k.textbox(s, x + 0.16, 1.82, 2.75, 0.7, title, 16, k.CYAN if sel else k.WHITE, True)
        k.textbox(s, x + 0.12, 2.60, 2.8, 2.8, body, 13, k.OFF)
    k.card(s, 0.42, 5.95, 12.48, 1.22, k.AMBER)
    k.textbox(s, 0.62, 6.08, 12.1, 0.28, "FROZEN CONFIGURATION", 11, k.AMBER, True, k.MONO)
    k.textbox(
        s, 0.62, 6.38, 12.1, 0.6,
        "4 rotors  ·  front pair tilts together  ·  rear pair tilts together  ·  max tilt 75°  ·  high wing + 2 winglets  ·  Pixhawk / ArduPilot",
        14, k.WHITE, True,
    )


def s09_architecture(prs):
    s = _base(
        prs, "08  /  ARCHITECTURE", "Overall Aircraft Architecture", 9,
        "Treat this as the product breakdown: airframe, propulsion, tilt, energy, GNC, payload.\n"
        "Front and rear tilt are mechanically identical — one design, two installs.\n"
        "Payload is a 1 kg mass dummy with a hard point on the CG.",
    )
    k.picture(s, A / "pushpak-sideview.png", 0.42, 1.22, 6.35, 3.55)
    k.hud_corners(s, 0.42, 1.22, 6.35, 3.55)
    blocks = [
        (0.42, 4.95, "AIRFRAME", "Fuselage 1.2–1.3 m\nWing 1.30 m · 0.32 m²\nWinglets × 2  ·  skids"),
        (3.62, 4.95, "PROPULSION", "4 × MY4215 400KV\n16-in class  ·  6S\nTarget 9 kgf static"),
        (6.82, 1.22, "TILT", "2 shafts  ·  2 servos\n8 bearings  ·  stops\nMag angle sensors"),
        (10.02, 1.22, "ENERGY", "6S Li-ion ~20 Ah\n1.0–1.5 kg pack\nCurrent / voltage log"),
        (6.82, 3.55, "GNC", "Pixhawk + ArduPilot\nESP32 companion\nPitot · GPS · IMU"),
        (10.02, 3.55, "PAYLOAD", "1.00 kg dummy\nCG hard point\nNot a sensor suite"),
        (6.82, 5.88, "FRONT AXIS", "Tilt together\nLeft + right motors"),
        (10.02, 5.88, "REAR AXIS", "Tilt together\nLeft + right motors"),
    ]
    for x, y, t, b in blocks:
        w, h = 3.05, 1.18
        if y < 4.9 and x >= 6.8:
            h = 2.15
        k.card(s, x, y, w, h, k.CYAN)
        k.textbox(s, x + 0.12, y + 0.08, w - 0.2, 0.28, t, 11, k.CYAN, True, k.MONO)
        k.multiline(s, x + 0.08, y + 0.36, w - 0.16, h - 0.42, b.split("\n"), 12, k.OFF)


def s10_propulsion(prs):
    s = _base(
        prs, "09  /  PROPULSION", "Propulsion System", 10,
        "MY4215 400KV on 6S is in the 16-inch efficiency band we need for hover margin without a 2 kg motor set.\n"
        "Four identical corners — no mixed motor types. That keeps the quad mixer honest in hover.\n"
        "Battery is Li-ion for energy density; C-rate will be confirmed on the thrust stand before first hover.",
    )
    kpis = [
        ("4 ×", "HGLRC MY4215", "400 KV  ·  200 g each"),
        ("800 g", "Total motor mass", "18% of 4.5 kg MTOW"),
        ("16 in", "Propeller class", "D = 406 mm"),
        ("6S", "Li-ion pack", "Initial concept ~20 Ah"),
    ]
    for i, (v, a, b) in enumerate(kpis):
        x = 0.42 + i * 3.20
        k.kpi(s, x, 1.22, 3.05, 1.55, v, a, b)

    data = [
        ["Item", "Spec", "Qty", "Role"],
        ["Motor", "HGLRC MY4215  400KV", "4", "Hover + cruise thrust"],
        ["Motor mass", "200 g  (800 g set)", "4", "Mass budget line"],
        ["Propeller", "16-inch class  /  406 mm", "4", "Disk area for hover"],
        ["ESC", "From 0.25 kg allocation", "4", "6S capable, logged current"],
        ["Battery", "6S Li-ion  ~20 Ah", "1", "Energy + CG"],
        ["Bus", "6S power distribution", "1", "Voltage / current sense"],
        ["Mounts", "Tilt-shaft motor mounts", "4", "Interface to tilt axis"],
    ]
    k.add_table(s, data, 0.42, 2.95, 12.48, 4.20, 12)
