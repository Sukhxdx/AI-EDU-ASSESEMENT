"""Slides 31–40 — test, results, risks, programme close."""

from pathlib import Path

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

import kit as k

A = Path(__file__).parent / "assets"


def _base(prs, kicker, title, num, notes, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    k.set_bg(s)
    k.header(s, kicker, title, subtitle)
    k.footer(s, num)
    k.notes(s, notes)
    return s


def s31_ground(prs):
    s = _base(
        prs, "30  /  GROUND TEST", "Ground Testing Plan", 31,
        "Nobody flies a 4.5 kg VTOL that has not been on a thrust stand and a tilt rig.\n"
        "GT-4 (tilt endurance) is the test that finds print-layer cracks.\n"
        "Power instrumentation is calibrated on the bench so the flight comparison is not a sensor argument.",
    )
    data = [
        ["ID", "Test", "Setup", "Pass"],
        ["GT-1", "Motor / prop map", "Single-motor stand, 6S", "≥ 2.25 kgf / motor; temps OK"],
        ["GT-2", "4-motor hover T", "Frame locked, pads", "≥ 9 kgf total; vibration log"],
        ["GT-3", "Tilt kinematics", "Axis jig, encoder", "0–75°, no bind, stop contact"],
        ["GT-4", "Tilt endurance", "500 cycles loaded", "No play, no mag-angle drift"],
        ["GT-5", "Servo stall / current", "Locked horn", "Current trip works"],
        ["GT-6", "Pitot calibration", "Wind / car boom", "IAS vs reference"],
        ["GT-7", "Power DAQ", "Known shunt load", "W within 5%"],
        ["GT-8", "EMI / compass", "Motors spinning", "Heading stable"],
        ["GT-9", "Failsafe dry-run", "FC on bench", "RTL, batt, GCS loss"],
        ["GT-10", "Tether hover", "Outdoor, short tether", "QHOVER 30 s, no oscillation"],
    ]
    k.add_table(s, data, 0.42, 1.18, 12.48, 5.95, 11)


def s32_flight(prs):
    s = _base(
        prs, "31  /  FLIGHT TEST", "Flight Testing Plan", 32,
        "Envelope expansion, not a show flight. FT-1 is a hover. Science happens at FT-6/7.\n"
        "Wing-off comparison is either a removable-wing article or a rotor-borne pass at the same mass with tilt locked at 0° in slow forward flight — we will pick one and freeze it in the test plan.\n"
        "Weather: low wind for first transitions.",
    )
    data = [
        ["Sortie", "Objective", "Mode", "Abort"],
        ["FT-1", "Hover 30–60 s, CG, vibration", "QHOVER", "Oscillation / mag-angle noise"],
        ["FT-2", "Pad ops, QLOITER box", "QLOITER", "GPS / compass disagree"],
        ["FT-3", "Low-speed translation, tilt 0°", "QHOVER", "Pilot PIO"],
        ["FT-4", "Partial tilt to 20–30°", "Transition", "Pitch upset, Vz band"],
        ["FT-5", "Full 75° to 20 m/s, short", "Trans. + FBWA", "Airspeed decay"],
        ["FT-6", "Wing-on power point 20 m/s", "CRUISE", "Current / temp"],
        ["FT-7", "Wing-off / rotor-borne compare", "VTOL fwd", "Power > limit"],
        ["FT-8", "Repeatability ×3", "FT-6 profile", "Inconsistent IAS"],
        ["FT-9", "QRTL from cruise", "QRTL", "Tilt reverse stall"],
        ["FT-10", "Demo for faculty", "Rehearsed profile", "Any red abort"],
    ]
    k.add_table(s, data, 0.42, 1.18, 12.48, 5.95, 11)


def s33_expected(prs):
    s = _base(
        prs, "32  /  RESULTS", "Expected Results", 33,
        "These are predictions, clearly labelled. Faculty should hear the uncertainty bands.\n"
        "The only result that defines success is a measured electrical-power reduction with uncertainty, at matched mass.\n"
        "CL 0.395 is a design point, not a trophy — we will report the α that actually produced it.",
    )
    kpis = [
        ("2.0×", "Hover T/W", "9 kgf static on 4.5 kg"),
        ("0.395", "Cruise CL", "70% weight at 20 m/s"),
        ("~40–55%", "Power cut (est.)", "Wing-on vs wing-off fwd"),
        ("≤ 3.5 kg", "Empty mass", "Hard engineering target"),
        ("75°", "Tilt stop", "Encoder-confirmed"),
        ("1.0 kg", "Payload", "Carried at design point"),
    ]
    for i, (v, a, b) in enumerate(kpis):
        x = 0.42 + (i % 3) * 4.16
        y = 1.22 + (i // 3) * 2.90
        k.card(s, x, y, 3.98, 2.70, k.CYAN)
        k.textbox(s, x + 0.18, y + 0.22, 3.62, 0.70, v, 28, k.WHITE, True, k.MONO)
        k.textbox(s, x + 0.18, y + 1.05, 3.62, 0.40, a.upper(), 14, k.CYAN, True)
        k.textbox(s, x + 0.18, y + 1.55, 3.62, 0.80, b, 14, k.OFF)


def s34_power(prs):
    s = _base(
        prs, "33  /  EXPERIMENT", "Wing-On vs Wing-Off Power Comparison", 34,
        "This is the money slide. Bars are momentum-theory plus a simple drag polar — not flight data.\n"
        "Wing-off forward at 20 m/s is the expensive condition: rotors still lifting, now also dragging.\n"
        "Wing-on at the 70% lift design point is conservative; a drag-sized thrust at 75° is the floor we hope to approach after tuning.\n"
        "Flight test will replace these bands with mean ± scatter from FT-6 and FT-7.",
        "Estimates only  ·  MTOW 4.5 kg  ·  6S  ·  16-inch class  ·  V = 20 m/s",
    )
    k.picture(s, A / "power_compare.png", 0.30, 1.18, 8.15, 4.55)
    k.hud_corners(s, 0.30, 1.18, 8.15, 4.55)
    k.card(s, 8.60, 1.18, 4.30, 4.55, k.CYAN)
    k.textbox(s, 8.78, 1.36, 4.0, 0.30, "WHAT WE WILL FLY", 12, k.CYAN, True, k.MONO)
    k.bullet_block(s, 8.70, 1.75, 4.05, 3.7, [
        "Same mass, same pack, same day.",
        "Stabilised 20 m/s run, level.",
        "Log W = V×I, IAS, tilt, RPM.",
        "Three repeats each config.",
        "Report mean, min, max — not a single heroic pass.",
        "If removable wing is unsafe, lock tilt at 0° for the rotor-borne leg.",
    ], 13)
    k.card(s, 0.42, 5.90, 12.48, 1.20, k.AMBER)
    k.textbox(
        s, 0.62, 6.15, 12.1, 0.75,
        "Success criterion   ·   statistically lower electrical power in the wing-on condition at matched MTOW and airspeed   ·   uncertainty reported",
        14, k.WHITE, True,
    )


def s35_risks(prs):
    s = _base(
        prs, "34  /  RISK", "Risks and Mitigation", 35,
        "Top risk is mass growth — it eats hover margin and the empty-mass promise.\n"
        "Second is transition pitch. Third is Li-ion C-rate on 16-inch props.\n"
        "We mitigate with a freeze date, a thrust stand, and a tether.",
    )
    data = [
        ["ID", "Risk", "L", "C", "Mitigation"],
        ["R1", "MTOW > 4.9 kg (battery / prints)", "M", "H", "Weekly CAD mass; foam wing; pack downsize"],
        ["R2", "Static thrust < 2.25 kgf / motor", "M", "H", "Stand map before CAD freeze; prop pitch options"],
        ["R3", "Transition pitch-up / ballooning", "M", "H", "CFD C4; slow tilt rate; abort to QHOVER"],
        ["R4", "Tilt jam / servo stall", "L", "H", "Stops, dual mag-angle, current trip, shaft bearings"],
        ["R5", "Li-ion sag under hover current", "M", "H", "Stand + IR test; keep LiPo backup pack"],
        ["R6", "Pitot / IAS false transition", "M", "M", "Fail-detect; inhibit tilt on IAS invalid"],
        ["R7", "EMI on compass / mag angle", "M", "M", "Routing, shielding, GT-8"],
        ["R8", "Schedule slip on manufacture", "H", "M", "COTS-first; parallel tilt bench article"],
        ["R9", "Loss of vehicle on first free flight", "M", "H", "Tether, geofence, two-person crew"],
        ["R10", "Inconclusive power comparison", "M", "H", "DAQ cal, repeats, locked test card"],
    ]
    k.add_table(s, data, 0.42, 1.18, 12.48, 5.95, 11)


def s36_future(prs):
    s = _base(
        prs, "35  /  NEXT", "Future Improvements", 36,
        "V1 is the experiment. V2 is the aeroplane we would build if V1's power delta is real.\n"
        "Do not promise autonomous delivery. Promise a cleaner tilt, a better polar, and maybe 90° if the data wants it.\n"
        "Differential front/rear tilt as a pitch effector is the most interesting control upgrade.",
    )
    items = [
        ("V1.1", "90° tilt option if 75° leaves unused axial efficiency."),
        ("V1.1", "Folding props or larger disk if hover thermal is the limiter."),
        ("V2", "Span / AR increase once manufacturing of the 1.3 m wing is proven."),
        ("V2", "Independent motor-pod tilt (true tilt-rotor) vs paired shafts."),
        ("V2", "Front/rear differential tilt as a pitch moment source."),
        ("V2", "Onboard real-time CL estimator from pitot + IMU + mass."),
        ("V3", "Redundant FC / dual battery if we ever leave the demo role."),
        ("V3", "Mission payload (EO) replacing the 1 kg dummy — only after power is shown."),
    ]
    for i, (tag, b) in enumerate(items):
        x = 0.42 + (i % 2) * 6.40
        y = 1.22 + (i // 2) * 1.40
        k.card(s, x, y, 6.20, 1.25, k.CYAN)
        k.textbox(s, x + 0.18, y + 0.16, 1.3, 0.90, tag, 16, k.CYAN, True, k.MONO, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        k.textbox(s, x + 1.55, y + 0.22, 4.45, 0.85, b, 14, k.OFF, False, k.FONT, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)


def s37_timeline(prs):
    s = _base(
        prs, "36  /  SCHEDULE", "Project Timeline", 37,
        "Eight months, one vehicle. The critical path is tilt CAD → printed axis → stand → tether.\n"
        "Flight test is not a week at the end — it is a month of envelope expansion.\n"
        "Faculty review gates sit at M2 (config freeze), M5 (prototype), M8 (results).",
    )
    k.picture(s, A / "timeline.png", 0.30, 1.18, 12.7, 4.15)
    gates = [
        ("M2", "PDR", "Config, mass, thrust frozen"),
        ("M5", "CDR / proto", "Vehicle complete, GT-10 done"),
        ("M7", "FT-6/7", "Power comparison in hand"),
        ("M8", "Closeout", "Report, demo, archive logs"),
    ]
    for i, (m, g, b) in enumerate(gates):
        x = 0.42 + i * 3.20
        k.card(s, x, 5.50, 3.05, 1.60, k.CYAN)
        k.textbox(s, x + 0.14, 5.62, 2.75, 0.26, m, 11, k.CYAN, True, k.MONO)
        k.textbox(s, x + 0.14, 5.90, 2.75, 0.32, g, 16, k.WHITE, True)
        k.textbox(s, x + 0.14, 6.28, 2.75, 0.60, b, 12, k.OFF)


def s38_budget(prs):
    s = _base(
        prs, "37  /  BUDGET", "Budget Estimate", 38,
        "Student-project money, not a company BOM. Figures are planning estimates in INR including contingency.\n"
        "Propulsion plus avionics are half the spend. The wing is cheap until we break it.\n"
        "Contingency 15% is for the second pack, the broken props, and the servo we will stall on the bench.",
        "Planning estimate  ·  COTS-heavy  ·  university workshop labour not costed",
    )
    data = [
        ["WBS", "Group", "Basis", "Est. (INR)"],
        ["1", "Propulsion (4× MY4215, props, ESCs)", "COTS", "28,000"],
        ["2", "Battery 6S Li-ion ~20 Ah + sense", "COTS", "22,000"],
        ["3", "Pixhawk + GPS/compass + pitot", "COTS", "32,000"],
        ["4", "ESP32, harness, PDB, connectors", "COTS", "6,000"],
        ["5", "Tilt (2× STS3215, bearings, shaft)", "COTS + shop", "12,000"],
        ["6", "Mag angle sensors ×2", "COTS", "5,000"],
        ["7", "Airframe (foam/composite, prints, gear)", "Shop", "18,000"],
        ["8", "Test (stand, shunt, prints, props spares)", "Shop", "8,000"],
        ["9", "Contingency 15%", "Risk", "19,650"],
        ["", "TOTAL (planning)", "", "≈ 1,50,650"],
    ]
    k.add_table(s, data, 0.42, 1.18, 12.48, 5.95, 12)


def s39_conclusion(prs):
    s = _base(
        prs, "38  /  CLOSE", "Conclusion", 39,
        "Restate the mission in one breath, then the frozen numbers, then the experiment.\n"
        "Pushpak V1 exists to measure a power delta, not to look like an eVTOL poster.\n"
        "Ask for the go-ahead to freeze configuration and buy long-lead motors, Pixhawk, and the 6S pack.",
    )
    k.card(s, 0.42, 1.22, 12.48, 1.15, k.CYAN)
    k.textbox(
        s, 0.62, 1.42, 12.1, 0.78,
        "Pushpak V1 is a 4.5 kg tilt-rotor VTOL demonstrator built to prove that wing-borne lift cuts electrical power versus rotor-borne flight.",
        16, k.WHITE, True,
    )
    points = [
        ("01", "Configuration frozen", "4 rotors, paired front/rear tilt to 75°, high wing NACA 4412, Pixhawk / ArduPilot / ESP32."),
        ("02", "Physics closed", "Hover 1.125 kgf/motor floor, 2.25 kgf design; cruise CL 0.395 at 20 m/s on 0.32 m²."),
        ("03", "Experiment defined", "Wing-on vs wing-off electrical power at matched mass — the only success metric that matters."),
        ("04", "Path to flight", "CAD → stand → tether → envelope expansion. Mass watched weekly against 3.5 kg empty."),
    ]
    for i, (n, t, b) in enumerate(points):
        y = 2.55 + i * 1.12
        k.card(s, 0.42, y, 12.48, 1.02, k.CYAN)
        k.textbox(s, 0.60, y + 0.28, 0.8, 0.48, n, 18, k.CYAN, True, k.MONO)
        k.textbox(s, 1.55, y + 0.12, 10.9, 0.32, t, 15, k.WHITE, True)
        k.textbox(s, 1.55, y + 0.48, 10.9, 0.42, b, 13, k.OFF)


def s40_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    k.set_bg(s)
    k.picture(s, A / "pushpak-hero.png", 5.35, 0, 7.983, 7.5)
    k.rect(s, 0, 0, 7.1, 7.5, k.BG)
    k.rect(s, 0, 0, 13.333, 0.055, k.CYAN)
    k.textbox(s, 0.50, 1.55, 6.3, 0.28, "PUSHPAK  V1", 12, k.CYAN, True, k.MONO)
    k.textbox(s, 0.48, 1.95, 6.3, 1.1, "Thank you", 44, k.WHITE, True)
    k.textbox(
        s, 0.50, 3.15, 6.1, 0.7,
        "Questions, technical discussion, and a configuration freeze.",
        16, k.OFF,
    )
    k.card(s, 0.50, 4.05, 6.05, 2.55, k.CYAN)
    k.textbox(s, 0.68, 4.20, 5.7, 0.28, "PROJECT TEAM", 11, k.CYAN, True, k.MONO)
    k.multiline(
        s, 0.62, 4.52, 5.75, 1.90,
        ["Ishaan Bafna      PRN 23070125017",
         "Aman Ravat        PRN 23070125004",
         "Krit Verma        PRN 230701250",
         "",
         "Tilt-Rotor VTOL UAV Demonstrator"],
        15, k.OFF, font=k.MONO,
    )
    k.footer(s, 40)
    k.notes(
        s,
        "Stop talking. Offer the test cards and the mass budget as handouts if asked.\n"
        "Likely questions: why 75° not 90°; why Li-ion; why paired shafts not four independent pods; how you will fly wing-off safely.\n"
        "Answers are on slides 8, 10, 16, and 34.",
    )
